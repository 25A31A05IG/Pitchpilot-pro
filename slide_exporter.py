import io
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_pitch_pptx(pitch_data: dict) -> io.BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    BG_COLOR = RGBColor(15, 23, 42)       # Slate 900
    ACCENT_COLOR = RGBColor(99, 102, 241) # Indigo 500
    TEXT_MAIN = RGBColor(248, 250, 252)   # Slate 50

    def apply_slide_bg(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    # Title Slide
    title_slide = prs.slides.add_slide(blank_layout)
    apply_slide_bg(title_slide)
    
    tb = title_slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(3.0))
    tf = tb.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = pitch_data.get("project_name", "PitchPilot Deck")
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = TEXT_MAIN

    p2 = tf.add_paragraph()
    p2.text = pitch_data.get("tagline", "")
    p2.font.size = Pt(22)
    p2.font.color.rgb = ACCENT_COLOR
    p2.space_before = Pt(14)

    # Content Slides
    for slide_info in pitch_data.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)
        apply_slide_bg(slide)

        h_box = slide.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.33), Inches(1.2))
        h_tf = h_box.text_frame
        h_p = h_tf.paragraphs[0]
        h_p.text = f"0{slide_info.get('slide_number', 1)} | {slide_info.get('title', '')}"
        h_p.font.size = Pt(28)
        h_p.font.bold = True
        h_p.font.color.rgb = ACCENT_COLOR

        b_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.33), Inches(4.5))
        b_tf = b_box.text_frame
        b_tf.word_wrap = True

        for bullet in slide_info.get("bullets", []):
            bp = b_tf.add_paragraph()
            bp.text = f"•  {bullet}"
            bp.font.size = Pt(20)
            bp.font.color.rgb = TEXT_MAIN
            bp.space_before = Pt(16)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output